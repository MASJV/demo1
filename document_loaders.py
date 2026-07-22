"""
TechTrinetra7 — Advanced RAG | Document Loaders
Supports: PDF, DOCX, TXT, MD, JSON, CSV, PPTX
All loaders are dependency-light and return LangChain `Document` objects
with rich metadata (source, page/slide/row, format) so later pipeline
stages (chunking, citation) can trace content back to its origin.
"""

from __future__ import annotations
import io
import json
import csv
import os
from typing import List

from langchain_core.documents import Document


# --------------------------------------------------------------------------
# Individual format loaders — each takes raw bytes + filename, returns Documents
# --------------------------------------------------------------------------

def load_pdf(file_bytes: bytes, filename: str) -> List[Document]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    docs = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            docs.append(
                Document(
                    page_content=text,
                    metadata={"source": filename, "format": "pdf", "page": i + 1},
                )
            )
    return docs


def load_docx(file_bytes: bytes, filename: str) -> List[Document]:
    import docx

    f = docx.Document(io.BytesIO(file_bytes))
    docs = []
    buffer, block_idx = [], 1

    def flush():
        nonlocal buffer, block_idx
        text = "\n".join(buffer).strip()
        if text:
            docs.append(
                Document(
                    page_content=text,
                    metadata={"source": filename, "format": "docx", "block": block_idx},
                )
            )
            block_idx += 1
        buffer = []

    for para in f.paragraphs:
        if para.style.name.startswith("Heading") and buffer:
            flush()
        if para.text.strip():
            buffer.append(para.text)
    flush()

    # Tables -> separate documents
    for t_idx, table in enumerate(f.tables):
        rows = []
        for row in table.rows:
            rows.append(" | ".join(c.text.strip() for c in row.cells))
        table_text = "\n".join(rows).strip()
        if table_text:
            docs.append(
                Document(
                    page_content=table_text,
                    metadata={"source": filename, "format": "docx", "table": t_idx + 1},
                )
            )
    return docs


def load_txt(file_bytes: bytes, filename: str) -> List[Document]:
    text = file_bytes.decode("utf-8", errors="ignore")
    return [Document(page_content=text, metadata={"source": filename, "format": "txt"})]


def load_md(file_bytes: bytes, filename: str) -> List[Document]:
    text = file_bytes.decode("utf-8", errors="ignore")
    return [Document(page_content=text, metadata={"source": filename, "format": "md"})]


def load_json(file_bytes: bytes, filename: str) -> List[Document]:
    data = json.loads(file_bytes.decode("utf-8", errors="ignore"))
    docs = []

    def stringify(obj, path="root"):
        if isinstance(obj, dict):
            for k, v in obj.items():
                stringify(v, f"{path}.{k}")
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                stringify(item, f"{path}[{i}]")
        else:
            docs.append(
                Document(
                    page_content=f"{path}: {obj}",
                    metadata={"source": filename, "format": "json", "path": path},
                )
            )

    # If it's a list of flat records (common case), keep each record as one chunk
    if isinstance(data, list) and all(isinstance(x, dict) for x in data):
        for i, record in enumerate(data):
            text = "\n".join(f"{k}: {v}" for k, v in record.items())
            docs.append(
                Document(
                    page_content=text,
                    metadata={"source": filename, "format": "json", "record": i},
                )
            )
    else:
        stringify(data)
        # collapse leaf-level documents into one readable blob for retrieval quality
        merged = "\n".join(d.page_content for d in docs)
        docs = [Document(page_content=merged, metadata={"source": filename, "format": "json"})]
    return docs


def load_csv(file_bytes: bytes, filename: str) -> List[Document]:
    import pandas as pd

    df = pd.read_csv(io.BytesIO(file_bytes))
    docs = []
    # Schema summary as its own document (helps retrieval for "what columns" queries)
    schema_text = f"CSV file '{filename}' columns: {', '.join(df.columns.astype(str))}. Rows: {len(df)}."
    docs.append(Document(page_content=schema_text, metadata={"source": filename, "format": "csv", "section": "schema"}))

    # Chunk rows in batches to avoid one-doc-per-row explosion on large files
    batch_size = 25
    for start in range(0, len(df), batch_size):
        batch = df.iloc[start:start + batch_size]
        text = batch.to_csv(index=False)
        docs.append(
            Document(
                page_content=text,
                metadata={"source": filename, "format": "csv", "rows": f"{start+1}-{start+len(batch)}"},
            )
        )
    return docs


def load_pptx(file_bytes: bytes, filename: str) -> List[Document]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    docs = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        texts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    texts.append(" | ".join(c.text for c in row.cells))
        slide_text = "\n".join(texts).strip()
        if slide_text:
            docs.append(
                Document(
                    page_content=slide_text,
                    metadata={"source": filename, "format": "pptx", "slide": i + 1},
                )
            )
    return docs


LOADER_MAP = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".txt": load_txt,
    ".md": load_md,
    ".json": load_json,
    ".csv": load_csv,
    ".pptx": load_pptx,
}


def load_any(file_bytes: bytes, filename: str) -> List[Document]:
    """Dispatch to the correct loader based on file extension."""
    ext = os.path.splitext(filename)[1].lower()
    if ext not in LOADER_MAP:
        raise ValueError(
            f"Unsupported file type '{ext}'. Supported: {', '.join(LOADER_MAP.keys())}"
        )
    return LOADER_MAP[ext](file_bytes, filename)


def load_many(files: List[tuple]) -> List[Document]:
    """files: list of (filename, bytes) tuples."""
    all_docs: List[Document] = []
    for filename, file_bytes in files:
        try:
            all_docs.extend(load_any(file_bytes, filename))
        except Exception as e:
            all_docs.append(
                Document(
                    page_content="",
                    metadata={"source": filename, "format": "error", "error": str(e)},
                )
            )
    return [d for d in all_docs if d.page_content.strip()]
